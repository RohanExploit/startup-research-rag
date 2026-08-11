from pathlib import Path

# Try to import necessary libraries for synthetic doc generation
try:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from faker import Faker
except ImportError:
    import subprocess
    import sys
    print("Installing required packages...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx", "openpyxl", "python-pptx", "reportlab", "faker"])
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from faker import Faker

fake = Faker()

def generate_pdf(path, title, paragraphs):
    c = canvas.Canvas(str(path))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 800, title)
    c.setFont("Helvetica", 12)
    y = 750
    for p in paragraphs:
        if y < 100:
            c.showPage()
            c.setFont("Helvetica", 12)
            y = 800
        # simple text wrapping (rudimentary)
        words = p.split()
        line = ""
        for word in words:
            if len(line) + len(word) > 80:
                c.drawString(72, y, line)
                y -= 20
                line = word + " "
            else:
                line += word + " "
        c.drawString(72, y, line)
        y -= 30
    c.save()

def generate_docx(path, title, paragraphs):
    doc = Document()
    doc.add_heading(title, 0)
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(path)

def generate_xlsx(path, title, rows):
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append([title])
    ws.append(["Name", "Department", "Salary", "Start Date"])
    for row in rows:
        ws.append(row)
    wb.save(path)

def generate_pptx(path, title, slides_content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    bullet_slide_layout = prs.slide_layouts[1]
    for slide_title, bullet_points in slides_content:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_title
        tf = body_shape.text_frame
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
    prs.save(path)


def main():
    out_dir = Path("data/tenants/tenant_1/raw")
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"Generating 50 synthetic documents in {out_dir}...")

    # 20 PDFs, 15 DOCX, 10 XLSX, 5 PPTX
    for i in range(20):
        path = out_dir / f"HR_Policy_{i+1:02d}.pdf"
        generate_pdf(path, f"HR Policy Document {i+1}", [fake.paragraph(nb_sentences=5) for _ in range(4)])

    for i in range(15):
        path = out_dir / f"Project_Proposal_{i+1:02d}.docx"
        generate_docx(path, f"Project Proposal {i+1}", [fake.paragraph(nb_sentences=5) for _ in range(3)])

    for i in range(10):
        path = out_dir / f"Financial_Report_{i+1:02d}.xlsx"
        rows = [[fake.name(), fake.job(), fake.random_int(min=50000, max=150000), str(fake.date_this_decade())] for _ in range(20)]
        generate_xlsx(path, f"Q{i%4 + 1} Financial Report", rows)

    for i in range(5):
        path = out_dir / f"Quarterly_Review_{i+1:02d}.pptx"
        slides = [
            ("Agenda", [fake.sentence(), fake.sentence()]),
            ("Key Metrics", [fake.sentence(), fake.sentence()]),
            ("Next Steps", [fake.sentence(), fake.sentence()])
        ]
        generate_pptx(path, f"Quarterly Review {i+1}", slides)

    print(f"Successfully generated 50 documents in {out_dir}")

if __name__ == "__main__":
    main()
